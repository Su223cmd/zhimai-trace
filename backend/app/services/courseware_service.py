import os
import uuid
import json
from datetime import datetime
from sqlalchemy.orm import Session
from pptx import Presentation
from pptx.util import Inches, Pt
import httpx

from app.models.db_models import Courseware, KnowledgePoint
from app.config import settings
from app.services.llm_client import call_llm_for_extraction, call_llm_for_relation_extraction, parse_json_response, _fallback_extraction

TEMPLATE_REGION_MARKERS = {
    "title": ["【标题】", "[标题]", "■ 标题"],
    "knowledge_points": ["【知识点】", "[知识点]", "■ 知识点"],
    "key_points": ["【重点】", "[重点]", "■ 重点"],
    "cognitive_level": ["【认知层级】", "[认知层级]", "■ 认知层级"],
}


def save_upload_file(upload_file, dest_dir: str) -> str:
    os.makedirs(dest_dir, exist_ok=True)
    file_id = str(uuid.uuid4())
    ext = os.path.splitext(upload_file.filename)[1] if upload_file.filename else ".pptx"
    filename = f"{file_id}{ext}"
    filepath = os.path.join(dest_dir, filename)
    with open(filepath, "wb") as f:
        content = upload_file.file.read()
        f.write(content)
    return filepath


def extract_ppt_content(ppt_path: str) -> list[dict]:
    prs = Presentation(ppt_path)
    slides = []
    for idx, slide in enumerate(prs.slides):
        slide_data = {
            "slide_index": idx + 1,
            "texts": [],
            "tables": [],
            "images": [],
            "notes": ""
        }
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_data["texts"].append(text)
            if shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                slide_data["tables"].append(table_data)
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            slide_data["notes"] = notes_text
        slides.append(slide_data)
    return slides


def detect_template_mode(slides: list[dict]) -> bool:
    required_markers = {"knowledge_points", "key_points"}
    matched_count = 0
    for slide in slides:
        all_text = " ".join(slide.get("texts", []))
        found = set()
        for region_name, markers in TEMPLATE_REGION_MARKERS.items():
            for marker in markers:
                if marker in all_text:
                    found.add(region_name)
                    break
        if required_markers.issubset(found):
            matched_count += 1
    return matched_count / len(slides) >= 0.8 if slides else False


def parse_template_slide(slide_data: dict) -> dict:
    lines = slide_data.get("texts", [])
    title = ""
    knowledge_points = []
    key_points = []
    cognitive_level = None
    current_region = None

    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue
        detected_region = None
        for region_name, markers in TEMPLATE_REGION_MARKERS.items():
            for marker in markers:
                if stripped.startswith(marker):
                    detected_region = region_name
                    stripped = stripped[len(marker):].strip()
                    break
            if detected_region:
                break

        if detected_region:
            current_region = detected_region

        if current_region == "title":
            title = stripped
        elif current_region == "knowledge_points":
            for item in stripped.replace("；", ";").split(";"):
                item = item.strip()
                if item:
                    knowledge_points.append(item)
        elif current_region == "key_points":
            for item in stripped.replace("；", ";").split(";"):
                item = item.strip()
                if item:
                    key_points.append(item)
        elif current_region == "cognitive_level":
            valid_levels = ["remember", "understand", "apply", "analyze", "evaluate", "create"]
            if stripped.lower() in valid_levels:
                cognitive_level = stripped.lower()

    return {
        "slide_index": slide_data["slide_index"],
        "title": title,
        "knowledge_points": knowledge_points,
        "key_points": key_points,
        "cognitive_level": cognitive_level,
        "parse_mode": "template",
        "needs_teacher_confirm": False
    }


async def parse_free_slide_with_llm(slide_data: dict) -> dict:
    texts = slide_data.get("texts", [])
    tables = slide_data.get("tables", [])
    title = texts[0] if texts else f"第{slide_data['slide_index']}页"

    slide_text = "\n".join(texts)
    if tables:
        for table in tables:
            for row in table:
                slide_text += "\n" + " | ".join(row)
    if slide_data.get("notes"):
        slide_text += "\n备注：" + slide_data["notes"]

    knowledge_points_raw = await call_llm_for_extraction(slide_text)

    knowledge_points = []
    key_points = []
    for kp in knowledge_points_raw:
        if isinstance(kp, dict):
            knowledge_points.append(kp.get("name", ""))
            if kp.get("is_key_point"):
                key_points.append(kp.get("name", ""))
        elif isinstance(kp, str):
            knowledge_points.append(kp)

    return {
        "slide_index": slide_data["slide_index"],
        "title": title,
        "knowledge_points": knowledge_points,
        "key_points": key_points,
        "cognitive_level": knowledge_points_raw[0].get("cognitive_level") if knowledge_points_raw and isinstance(knowledge_points_raw[0], dict) else None,
        "parse_mode": "free",
        "needs_teacher_confirm": True,
        "slide_text": slide_text,
    }


async def parse_courseware(courseware_id: str, db: Session) -> dict:
    courseware = db.query(Courseware).filter(Courseware.id == courseware_id).first()
    if not courseware:
        raise ValueError(f"Courseware {courseware_id} not found")
    if not os.path.exists(courseware.file_path):
        raise ValueError(f"File not found: {courseware.file_path}")

    slides = extract_ppt_content(courseware.file_path)
    is_template = detect_template_mode(slides)
    parse_mode = "template" if is_template else "free"

    parsed_slides = []
    for slide in slides:
        if is_template:
            parsed = parse_template_slide(slide)
        else:
            parsed = await parse_free_slide_with_llm(slide)
        parsed_slides.append(parsed)

    created_kps = []
    for slide_parsed in parsed_slides:
        for idx, kp_name in enumerate(slide_parsed.get("knowledge_points", [])):
            code = f"{courseware.subject.upper()}-CW{str(courseware.id)[:8]}-S{slide_parsed['slide_index']}-KP{idx+1:02d}"
            kp = KnowledgePoint(
                code=code,
                name=kp_name,
                cognitive_level=slide_parsed.get("cognitive_level"),
                chapter=slide_parsed.get("title", ""),
                source_type=parse_mode,
                courseware_id=courseware.id,
            )
            db.add(kp)
            created_kps.append({
                "code": code,
                "name": kp_name,
                "slide_index": slide_parsed["slide_index"],
                "cognitive_level": slide_parsed.get("cognitive_level"),
            })

    courseware.parse_status = "completed"
    courseware.parse_mode = parse_mode
    courseware.slide_count = len(slides)
    courseware.parsed_at = datetime.now()
    db.commit()

    sync_courseware_knowledge({
        "courseware_id": str(courseware.id),
        "name": courseware.name,
        "parse_mode": parse_mode,
        "parsed_slides": parsed_slides,
    }, db)

    edc_result = await edc_extract_from_courseware_async(parsed_slides)

    persist_to_db(db)

    from app.services.agent_bus import AgentBus
    AgentBus.update_agent_state(db, "knowledge", "courseware_synced", {"last_courseware_id": str(courseware.id)})
    AgentBus.send(db, "knowledge", "diagnosis", "graph_updated", {
        "source": "courseware",
        "courseware_id": str(courseware.id),
        "kp_count": len(created_kps),
    })
    AgentBus.notify_teacher(db, 1, "knowledge_sync", "知识图谱已更新",
                            f"课件「{courseware.name}」已同步，新增{len(created_kps)}个知识点")

    return {
        "courseware_id": str(courseware.id),
        "parse_mode": parse_mode,
        "total_slides": len(slides),
        "parsed_slides": parsed_slides,
        "created_knowledge_points": len(created_kps),
        "knowledge_points": created_kps,
        "edc_result": edc_result,
    }


async def edc_extract_from_courseware_async(parsed_slides: list[dict]) -> dict:
    from app.services.knowledge_service import edc_pipeline, add_relation

    all_text = ""
    entity_names = []
    for slide in parsed_slides:
        kps = slide.get("knowledge_points", [])
        if isinstance(kps, list):
            for t in kps:
                name = t if isinstance(t, str) else t.get("name", "")
                if name:
                    all_text += name + "。"
                    entity_names.append(name)
        title = slide.get("title", "")
        if title:
            all_text += title + "。"
        slide_text = slide.get("slide_text", "")
        if slide_text:
            all_text += slide_text + "。"

    pattern_result = edc_pipeline(all_text) if all_text.strip() else {
        "total_extracted": 0, "defined": 0, "canonicalized": 0,
        "pending_review": 0, "confirmed_relations": [], "pending_relations": [],
    }

    llm_relations = []
    if entity_names and settings.deepseek_api_key:
        llm_extracted = await call_llm_for_relation_extraction(all_text[:3000], entity_names)
        for rel in llm_extracted:
            from app.services.knowledge_service import edc_define, edc_canonicalize, _get_store
            defined = edc_define([rel])
            kp_entities = _get_store()["entities"].get("KnowledgePoint", {})
            canonicalized = edc_canonicalize(defined, kp_entities)
            for t in canonicalized:
                if t.get("canonicalized"):
                    add_relation(
                        t["source_code"], t["target_code"], t["relation"],
                        weight=t.get("confidence", 0.5),
                        confidence=t.get("confidence", 0.5),
                        discovered_by="edc_llm",
                        evidence={"method": "llm", "source_text": t.get("source_text", "")},
                    )
                    llm_relations.append(t)

    pattern_confirmed = pattern_result.get("confirmed_relations", [])
    return {
        "total_extracted": pattern_result.get("total_extracted", 0) + len(llm_relations),
        "defined": pattern_result.get("defined", 0) + len(llm_relations),
        "canonicalized": pattern_result.get("canonicalized", 0) + len([r for r in llm_relations if r.get("canonicalized")]),
        "pending_review": pattern_result.get("pending_review", 0),
        "confirmed_relations": pattern_confirmed + llm_relations,
        "pending_relations": pattern_result.get("pending_relations", []),
        "pattern_extracted": pattern_result.get("total_extracted", 0),
        "llm_extracted": len(llm_relations),
    }


def create_template_ppt(output_path: str) -> str:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_layout = prs.slide_layouts[1]

    example_slides = [
        {
            "title": "大气的受热过程",
            "knowledge_points": "大气受热过程；大气逆辐射；温室效应",
            "key_points": "太阳辐射是根本能源；地面是直接热源；大气逆辐射起保温作用",
            "cognitive_level": "understand"
        },
        {
            "title": "热力环流",
            "knowledge_points": "热力环流原理；等压面弯曲；山谷风",
            "key_points": "冷热不均是根本原因；气压梯度力是直接原因",
            "cognitive_level": "apply"
        },
    ]

    for example in example_slides:
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = f"【标题】{example['title']}"

        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()

        paragraphs = [
            f"【知识点】{example['knowledge_points']}",
            f"【重点】{example['key_points']}",
            f"【认知层级】{example['cognitive_level']}",
        ]
        for i, text in enumerate(paragraphs):
            if i == 0:
                tf.paragraphs[0].text = text
            else:
                p = tf.add_paragraph()
                p.text = text
                p.font.size = Pt(18)

    prs.save(output_path)
    return output_path
